"""Parser for PowerPoint (.pptx) files — direct text-level PII sanitisation."""

import io
from pptx import Presentation
from app.pii_engine import sanitize_batch


def parse_pptx(data: bytes, mode: str = "redact") -> tuple[bytes, int]:
    prs = Presentation(io.BytesIO(data))

    # Collect all paragraph texts and their references
    all_texts = []
    para_refs = []  # store (paragraph, original_text) for write-back

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        all_texts.append(para.text)
                        para_refs.append(para)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            if para.text.strip():
                                all_texts.append(para.text)
                                para_refs.append(para)

    if not all_texts:
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue(), 0

    # Single batch sanitize
    cleaned_texts, total = sanitize_batch(all_texts, mode)

    for i, para in enumerate(para_refs):
        if cleaned_texts[i] != all_texts[i] and para.runs:
            para.runs[0].text = cleaned_texts[i]
            for run in para.runs[1:]:
                run.text = ""

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), total
